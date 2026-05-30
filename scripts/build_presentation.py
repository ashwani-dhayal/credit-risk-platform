"""Build the solution deck as both a PowerPoint and a PDF.

Outputs (under documents/):
  - documents/presentation.pptx   (editable source)
  - documents/presentation.pdf    (frozen, committed)

The deck mixes:
  * narrative slides (text bullets, metrics table)
  * frontend screenshots from documents/screenshots/
  * backend code snippets pulled live from src/

This means re-running this script always produces an in-sync deck.
"""

from __future__ import annotations

import json
import sys
from pathlib import Path
from typing import Iterable

PROJECT_ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(PROJECT_ROOT))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu, Inches, Pt
from PIL import Image  # ships with reportlab dependency tree
from reportlab.lib import colors
from reportlab.lib.pagesizes import landscape, A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import cm
from reportlab.platypus import (
    Image as RLImage,
    PageBreak,
    Paragraph,
    Preformatted,
    SimpleDocTemplate,
    Spacer,
    Table,
    TableStyle,
)


METRICS_PATH = PROJECT_ROOT / "models" / "metrics.json"
RULES_PATH = PROJECT_ROOT / "models" / "rules.json"
SCREENSHOT_DIR = PROJECT_ROOT / "documents" / "screenshots"
SRC_DIR = PROJECT_ROOT / "src"


def _load_metrics() -> dict:
    if METRICS_PATH.exists():
        return json.loads(METRICS_PATH.read_text())
    return {}


def _load_rules() -> list[dict]:
    if RULES_PATH.exists():
        return json.loads(RULES_PATH.read_text())
    return []


def _read_snippet(rel_path: str, start: int, end: int) -> str:
    """Return lines [start..end] (1-indexed, inclusive) from a project file."""
    full = PROJECT_ROOT / rel_path
    if not full.exists():
        return f"# {rel_path} not found"
    lines = full.read_text(encoding="utf-8").splitlines()
    chunk = lines[max(0, start - 1):end]
    return "\n".join(chunk)


# ----------------------------- Slide spec -----------------------------------
def _slide_data() -> list[dict]:
    m = _load_metrics()
    rules = _load_rules()[:5]

    slides: list[dict] = [
        # ---- 1. Title ----
        {
            "kind": "title",
            "title": "AI-Powered Credit Risk Intelligence Platform",
            "subtitle": "NeoStats AI Engineer Assignment — Submission",
            "bullets": [
                "Dataset: Home Credit Default Risk (Kaggle)",
                "Stack: Python 3.11, LightGBM, SHAP, Streamlit, SQLite, Docker",
                "LLM: Multi-provider (OpenAI / Groq / Gemini) + deterministic fallback",
                "Author: Ashwani Dhayal",
            ],
        },
        # ---- 2. Problem ----
        {
            "kind": "bullets",
            "title": "Business Problem",
            "bullets": [
                "Banks must make faster, more accurate, and explainable credit decisions.",
                "Identify high-risk applicants early and automate risk scoring.",
                "Provide auditable, regulator-friendly reasons for every decision.",
                "Let business analysts explore the portfolio in plain English.",
                "Bridge ML insights and credit policy through readable rules.",
            ],
        },
        # ---- 3. Solution at a glance ----
        {
            "kind": "bullets",
            "title": "Solution at a Glance",
            "bullets": [
                "End-to-end platform: EDA → ML → Explainability → Rules → Chatbot.",
                "Single Streamlit UI with 6 sections; one Docker command to run.",
                "ML model returns probability + Low/Medium/High band + Decision.",
                "SHAP TreeExplainer attributes every score to ranked feature drivers.",
                "Talk-to-Data agent translates English to safe, read-only SQLite SQL.",
                "Decision-tree rules give credit officers an audit-ready policy view.",
            ],
        },
        # ---- 4. Architecture ----
        {
            "kind": "bullets",
            "title": "Architecture (logical view)",
            "bullets": [
                "Streamlit UI (app/streamlit_app.py) — multi-section navigation.",
                "src/ml: LightGBM training / inference / SHAP TreeExplainer.",
                "src/rules: depth-4 decision tree → human-readable IF-THEN rules.",
                "src/data: schema, sample generator, CSV → SQLite ingestion.",
                "src/llm: OpenAI / Groq / Gemini auto-detect + safe NL→SQL agent.",
                "src/utils/sql_safety: read-only SELECT guardrails (no DDL/DML).",
                "Single SQLite DB is the source of truth for ML and the chatbot.",
            ],
        },
        # ---- 5. FRONTEND: Overview screenshot ----
        {
            "kind": "screenshot",
            "title": "UI · Overview Section",
            "image": str(SCREENSHOT_DIR / "01_overview.png"),
            "caption": (
                "Live KPIs of the platform: 10,000 applicants, 8% default rate, "
                "28 features after engineering, model ROC-AUC 0.895. The sidebar "
                "shows the active LLM provider and the loaded model artifact."
            ),
        },
        # ---- 6. FRONTEND: EDA screenshot ----
        {
            "kind": "screenshot",
            "title": "UI · Exploratory Data Analysis",
            "image": str(SCREENSHOT_DIR / "02_eda.png"),
            "caption": (
                "Four EDA tabs: Summary, Demographics, Financials, Default Drivers. "
                "Distributions are split by default status; the table on the right "
                "shows missingness for the most-affected columns."
            ),
        },
        # ---- 7. FRONTEND: Predict screenshot ----
        {
            "kind": "screenshot",
            "title": "UI · Risk Prediction",
            "image": str(SCREENSHOT_DIR / "03_predict.png"),
            "caption": (
                "Form-driven applicant scoring. Output: probability of default + "
                "risk band (Low / Medium / High) + suggested decision "
                "(Approve / Review / Reject). Sample shown: P(default)=0.06% → Low → Approve."
            ),
        },
        # ---- 8. FRONTEND: Explainability screenshot ----
        {
            "kind": "screenshot",
            "title": "UI · Explainability (SHAP)",
            "image": str(SCREENSHOT_DIR / "04_explain.png"),
            "caption": (
                "Per-applicant top SHAP contributors. Bars going right increase "
                "predicted default probability; bars going left decrease it. The "
                "table below shows the exact feature values and SHAP magnitudes."
            ),
        },
        # ---- 9. FRONTEND: Rules screenshot ----
        {
            "kind": "screenshot",
            "title": "UI · Decision Rules",
            "image": str(SCREENSHOT_DIR / "05_rules.png"),
            "caption": (
                "Business-readable IF-THEN rules derived from a depth-4 decision "
                "tree. Each row reports support (% of population), default rate "
                "inside the leaf, and lift vs. the 8% base rate."
            ),
        },
        # ---- 10. FRONTEND: Chatbot screenshot ----
        {
            "kind": "screenshot",
            "title": "UI · Talk-to-Data Chatbot",
            "image": str(SCREENSHOT_DIR / "06_chatbot.png"),
            "caption": (
                "User asks a plain-English question; the agent emits a safe "
                "SELECT, runs it on SQLite, and summarises the rows. The "
                "generated SQL is exposed for full transparency. Falls back to "
                "a deterministic intent parser when no LLM key is set."
            ),
        },
        # ---- 11. BACKEND: ML training code snippet ----
        {
            "kind": "code",
            "title": "Backend · ML Training (src/ml/train.py)",
            "subtitle": (
                "LightGBM with scale_pos_weight for class imbalance, threshold "
                "chosen by Youden's J on the validation set."
            ),
            "code": _read_snippet("src/ml/train.py", 70, 110),
        },
        # ---- 12. BACKEND: NL→SQL pipeline ----
        {
            "kind": "code",
            "title": "Backend · NL → SQL Agent (src/llm/nl_to_sql.py)",
            "subtitle": (
                "Pipeline: prompt LLM → parse JSON → static safety check → "
                "read-only execute → second LLM call grounds the summary in rows."
            ),
            "code": _read_snippet("src/llm/nl_to_sql.py", 178, 230),
        },
        # ---- 13. BACKEND: SQL safety guardrails ----
        {
            "kind": "code",
            "title": "Backend · SQL Safety (src/utils/sql_safety.py)",
            "subtitle": (
                "Single SELECT, table allowlist, blocked keyword list, automatic "
                "LIMIT cap. The chatbot cannot mutate data even if the LLM goes off-script."
            ),
            "code": _read_snippet("src/utils/sql_safety.py", 18, 70),
        },
        # ---- 14. BACKEND: Multi-provider LLM client ----
        {
            "kind": "code",
            "title": "Backend · Multi-Provider LLM Client (src/llm/client.py)",
            "subtitle": (
                "Auto-detects OpenAI / Groq / Gemini based on whichever API "
                "key is set; gracefully falls back to a deterministic parser."
            ),
            "code": _read_snippet("src/llm/client.py", 28, 86),
        },
        # ---- 15. Token optimisation ----
        {
            "kind": "bullets",
            "title": "Prompt Engineering & Token Optimisation",
            "bullets": [
                "JSON-only output contract — model emits {\"sql\": \"...\"}; deterministic parsing.",
                "Pinned schema (~25 cols, 1-line descriptions) → no DB introspection round-trip.",
                "Few-shot kept tiny (4 examples) to anchor style without bloating context.",
                "Temperature 0; max 256 output tokens — SQL doesn't need prose.",
                "Static SQL safety: SELECT/WITH only, table allowlist, LIMIT cap, blocked keywords.",
                "Result-grounded summary prompt explicitly forbids invented numbers.",
                "Hard fallback: deterministic regex parser keeps the demo working offline.",
            ],
        },
        # ---- 16. Evaluation results table ----
        {
            "kind": "metrics_table",
            "title": "Evaluation Results (bundled 10k sample)",
            "metrics_table": [
                ["Metric", "Value"],
                ["ROC-AUC", str(m.get("roc_auc", "—"))],
                ["PR-AUC", str(m.get("pr_auc", "—"))],
                ["KS Statistic", str(m.get("ks_statistic", "—"))],
                ["F1 @ optimal threshold", str(m.get("f1_at_optimal_threshold", "—"))],
                ["Optimal threshold", str(m.get("optimal_threshold", "—"))],
                ["Default rate", str(m.get("default_rate", "—"))],
                ["Train rows / Val rows", f"{m.get('n_train','—')} / {m.get('n_val','—')}"],
            ],
        },
        # ---- 17. Top derived rules ----
        {
            "kind": "bullets",
            "title": "Top Derived Decision Rules",
            "bullets": [
                f"R{r['rule_id']} · {r['band']} · support {r['support_pct']}% · "
                f"default rate {r['default_rate_pct']}% · lift {r['lift']}×"
                for r in rules
            ] or ["(rules.json not generated yet)"],
        },
        # ---- 18. How to run ----
        {
            "kind": "bullets",
            "title": "How To Run",
            "bullets": [
                "git clone https://github.com/ashwani-dhayal/credit-risk-platform.git",
                "cd credit-risk-platform",
                "cp .env.example .env  # add ONE LLM key (optional)",
                "docker compose up --build",
                "→ open http://localhost:8501",
                "Local: python -m venv .venv ; pip install -r requirements.txt ; "
                "python scripts/train_model.py ; streamlit run app/streamlit_app.py",
            ],
        },
        # ---- 19. Limitations ----
        {
            "kind": "bullets",
            "title": "Limitations & Next Steps",
            "bullets": [
                "Adding bureau / previous_application aggregates would lift AUC ~+0.03.",
                "No data-drift monitoring yet (Evidently / WhyLogs are obvious next steps).",
                "Single-tenant SQLite — switch to Postgres for multi-user concurrency.",
                "No UI auth — deploy behind a reverse proxy with auth in production.",
                "LLM cost cap is per-call only; daily quotas would need Redis-based limiting.",
            ],
        },
        # ---- 20. Thank you ----
        {
            "kind": "bullets",
            "title": "Thank You",
            "bullets": [
                "Repo: https://github.com/ashwani-dhayal/credit-risk-platform",
                "Author: Ashwani Dhayal",
                "Submission: NeoStats AI Engineer assignment.",
            ],
        },
    ]
    return slides


# ============================================================================
# PowerPoint
# ============================================================================
PRIMARY = RGBColor(0x1F, 0x77, 0xB4)


def _set_title(slide, text: str) -> None:
    if slide.shapes.title is None:
        return
    slide.shapes.title.text = text
    for para in slide.shapes.title.text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = PRIMARY
            run.font.bold = True


def _add_bullets(slide, bullets: Iterable[str]) -> None:
    placeholder = None
    for shp in slide.placeholders:
        if shp.placeholder_format.idx == 1:
            placeholder = shp
            break
    if placeholder is None:
        return
    tf = placeholder.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        for run in p.runs:
            run.font.size = Pt(18)


def _fit_image(image_path: Path, slide_w: int, slide_h: int) -> tuple[int, int, int, int]:
    """Compute centred (left, top, width, height) so the image fits the slide."""
    with Image.open(image_path) as im:
        iw, ih = im.size
    target_w = slide_w - Inches(0.6) * 2
    target_h = slide_h - Inches(1.6) - Inches(1.0)  # leave room for title + caption
    scale = min(target_w / iw, target_h / ih)
    w = int(iw * scale)
    h = int(ih * scale)
    left = int((slide_w - w) / 2)
    top = Inches(1.6)
    return left, int(top), w, h


def _build_pptx(slides: list[dict], path: Path) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title_layout = prs.slide_layouts[0]
    bullet_layout = prs.slide_layouts[1]
    blank_layout = prs.slide_layouts[5]

    for i, s in enumerate(slides):
        kind = s["kind"]

        if kind == "title":
            slide = prs.slides.add_slide(title_layout)
            _set_title(slide, s["title"])
            slide.placeholders[1].text = s.get("subtitle", "")
            # Append the bullet list under the subtitle
            box = slide.shapes.add_textbox(Inches(1.0), Inches(4.0),
                                           Inches(11.5), Inches(2.6))
            tf = box.text_frame
            tf.word_wrap = True
            for j, b in enumerate(s.get("bullets", [])):
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                p.text = "• " + b
                for run in p.runs:
                    run.font.size = Pt(16)

        elif kind == "bullets":
            slide = prs.slides.add_slide(bullet_layout)
            _set_title(slide, s["title"])
            _add_bullets(slide, s["bullets"])

        elif kind == "metrics_table":
            slide = prs.slides.add_slide(blank_layout)
            _set_title(slide, s["title"])
            rows = s["metrics_table"]
            tbl = slide.shapes.add_table(
                rows=len(rows), cols=2,
                left=Inches(2.5), top=Inches(1.7),
                width=Inches(8.3), height=Inches(0.45 * len(rows)),
            ).table
            for r_idx, row in enumerate(rows):
                for c_idx, val in enumerate(row):
                    cell = tbl.cell(r_idx, c_idx)
                    cell.text = str(val)
                    for para in cell.text_frame.paragraphs:
                        for run in para.runs:
                            run.font.size = Pt(16)
                            if r_idx == 0:
                                run.font.bold = True

        elif kind == "screenshot":
            slide = prs.slides.add_slide(blank_layout)
            _set_title(slide, s["title"])
            img_path = Path(s["image"])
            if img_path.exists():
                left, top, w, h = _fit_image(img_path, prs.slide_width, prs.slide_height)
                slide.shapes.add_picture(str(img_path), left, top, width=w, height=h)
            cap = slide.shapes.add_textbox(
                Inches(0.5), prs.slide_height - Inches(1.0),
                prs.slide_width - Inches(1.0), Inches(0.9),
            )
            ctf = cap.text_frame
            ctf.word_wrap = True
            ctf.text = s.get("caption", "")
            for para in ctf.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(13)
                    run.font.italic = True

        elif kind == "code":
            slide = prs.slides.add_slide(blank_layout)
            _set_title(slide, s["title"])
            sub = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.4),
                prs.slide_width - Inches(1.0), Inches(0.6),
            )
            stf = sub.text_frame
            stf.word_wrap = True
            stf.text = s.get("subtitle", "")
            for para in stf.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(13)
                    run.font.italic = True
                    run.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

            box = slide.shapes.add_textbox(
                Inches(0.5), Inches(2.0),
                prs.slide_width - Inches(1.0), Inches(5.2),
            )
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(0xF6, 0xF8, 0xFA)
            tf = box.text_frame
            tf.word_wrap = True
            tf.margin_left = Pt(8)
            tf.margin_right = Pt(8)
            tf.margin_top = Pt(8)
            tf.margin_bottom = Pt(8)
            tf.text = s["code"]
            for para in tf.paragraphs:
                for run in para.runs:
                    run.font.name = "Consolas"
                    run.font.size = Pt(11)

    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(path)
    return path


# ============================================================================
# PDF (ReportLab)
# ============================================================================
def _build_pdf(slides: list[dict], path: Path) -> Path:
    path.parent.mkdir(parents=True, exist_ok=True)
    page = landscape(A4)
    page_w, page_h = page

    doc = SimpleDocTemplate(
        str(path),
        pagesize=page,
        leftMargin=1.4 * cm,
        rightMargin=1.4 * cm,
        topMargin=1.0 * cm,
        bottomMargin=1.0 * cm,
        title="Credit Risk Intelligence Platform",
        author="Ashwani Dhayal",
    )

    styles = getSampleStyleSheet()
    h1 = ParagraphStyle(
        "H1", parent=styles["Heading1"],
        fontSize=22, leading=26,
        textColor=colors.HexColor("#1f77b4"),
        spaceAfter=10,
    )
    sub = ParagraphStyle(
        "Sub", parent=styles["Heading3"],
        fontSize=13, leading=17,
        textColor=colors.HexColor("#555555"),
        spaceAfter=10,
    )
    body = ParagraphStyle(
        "Body", parent=styles["BodyText"],
        fontSize=12, leading=17, spaceAfter=4,
    )
    cap = ParagraphStyle(
        "Caption", parent=styles["Italic"],
        fontSize=10, leading=14, textColor=colors.HexColor("#444444"),
        spaceBefore=6,
    )
    code = ParagraphStyle(
        "Code", parent=styles["Code"],
        fontSize=8.5, leading=11,
        backColor=colors.HexColor("#f6f8fa"),
        borderPadding=4,
    )

    flow: list = []
    inner_w = page_w - 2.8 * cm

    for i, s in enumerate(slides):
        kind = s["kind"]
        flow.append(Paragraph(s["title"], h1))

        if kind == "title":
            if s.get("subtitle"):
                flow.append(Paragraph(s["subtitle"], sub))
            for b in s.get("bullets", []):
                flow.append(Paragraph(f"• {b}", body))

        elif kind == "bullets":
            for b in s.get("bullets", []):
                flow.append(Paragraph(f"• {b}", body))

        elif kind == "metrics_table":
            tbl = Table(s["metrics_table"], colWidths=[8 * cm, 8 * cm])
            tbl.setStyle(TableStyle([
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#1f77b4")),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                ("FONTSIZE", (0, 0), (-1, -1), 12),
                ("BOTTOMPADDING", (0, 0), (-1, 0), 8),
                ("BACKGROUND", (0, 1), (-1, -1), colors.HexColor("#f6f8fa")),
                ("GRID", (0, 0), (-1, -1), 0.4, colors.lightgrey),
                ("ALIGN", (0, 0), (-1, -1), "LEFT"),
                ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                ("LEFTPADDING", (0, 0), (-1, -1), 10),
            ]))
            flow.append(tbl)

        elif kind == "screenshot":
            img_path = Path(s["image"])
            if img_path.exists():
                with Image.open(img_path) as im:
                    iw, ih = im.size
                target_h = page_h - 5.0 * cm
                target_w = inner_w
                scale = min(target_w / iw, target_h / ih)
                w = iw * scale
                h = ih * scale
                flow.append(RLImage(str(img_path), width=w, height=h))
            if s.get("caption"):
                flow.append(Paragraph(s["caption"], cap))

        elif kind == "code":
            if s.get("subtitle"):
                flow.append(Paragraph(s["subtitle"], sub))
            flow.append(Preformatted(s["code"], code))

        if i != len(slides) - 1:
            flow.append(PageBreak())

    doc.build(flow)
    return path


def main() -> int:
    slides = _slide_data()
    out_dir = PROJECT_ROOT / "documents"
    pptx_path = _build_pptx(slides, out_dir / "presentation.pptx")
    pdf_path = _build_pdf(slides, out_dir / "presentation.pdf")
    print(f">> PPTX -> {pptx_path}  ({len(slides)} slides)")
    print(f">> PDF  -> {pdf_path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
